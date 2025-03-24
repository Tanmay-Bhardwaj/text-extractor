from flask import Flask, request, jsonify, send_file
import os
import tempfile
from werkzeug.utils import secure_filename
import io

# Document processing libraries
try:
    import docx
    import pptx
    import PyPDF2
    import markdown
except ImportError:
    print("Please install required packages using:")
    print("pip install python-docx python-pptx PyPDF2 markdown flask")

app = Flask(__name__)

# Configure a temporary folder to store uploaded files
UPLOAD_FOLDER = tempfile.gettempdir()
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Allowed file extensions
ALLOWED_EXTENSIONS = {'txt', 'docx', 'pptx', 'pdf', 'md'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_from_txt(file_path):
    """Extract text from a .txt file"""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
        return file.read()

def extract_from_docx(file_path):
    """Extract text from a .docx file"""
    doc = docx.Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def extract_from_pptx(file_path):
    """Extract text from a .pptx file"""
    presentation = pptx.Presentation(file_path)
    text = []
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return '\n'.join(text)

def extract_from_pdf(file_path):
    """Extract text from a .pdf file"""
    reader = PyPDF2.PdfReader(file_path)
    text = []
    
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text.append(page.extract_text())
    
    return '\n'.join(text)

def extract_from_md(file_path):
    """Extract text from a .md file"""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
        md_content = file.read()
        # Convert markdown to HTML, then strip HTML tags for plain text
        html = markdown.markdown(md_content)
        # A simple way to strip HTML tags (for basic usage)
        import re
        text = re.sub(r'<[^>]+>', '', html)
        return text

def extract_text(file_path):
    """Extract text from a file based on its extension"""
    extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if extension == '.txt':
            return extract_from_txt(file_path)
        elif extension == '.docx':
            return extract_from_docx(file_path)
        elif extension == '.pptx':
            return extract_from_pptx(file_path)
        elif extension == '.pdf':
            return extract_from_pdf(file_path)
        elif extension == '.md':
            return extract_from_md(file_path)
        else:
            return f"Error: Unsupported file extension '{extension}'."
    except Exception as e:
        return f"Error extracting text from '{file_path}': {str(e)}"

@app.route('/')
def index():
    return app.send_static_file('index.html')

@app.route('/extract', methods=['POST'])
def extract():
    # Check if the post request has the file part
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    
    # If user does not select file, browser also
    # submits an empty part without filename
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text from the file
        extracted_text = extract_text(file_path)
        
        # Clean up - remove the file
        os.remove(file_path)
        
        # Return the extracted text
        return jsonify({'text': extracted_text})
    
    return jsonify({'error': 'File type not allowed'}), 400

@app.route('/download', methods=['POST'])
def download():
    data = request.json
    if not data or 'text' not in data:
        return jsonify({'error': 'No text provided'}), 400
    
    text = data['text']
    filename = data.get('filename', 'extracted_text.txt')
    
    # Create a file-like object from the text
    text_io = io.BytesIO(text.encode('utf-8'))
    
    # Return the file
    return send_file(
        text_io,
        as_attachment=True,
        download_name=filename,
        mimetype='text/plain'
    )

if __name__ == '__main__':
    # Create static folder if it doesn't exist
    if not os.path.exists('static'):
        os.mkdir('static')
    
    # Copy the HTML file to the static folder
    with open('static/index.html', 'w', encoding='utf-8') as f:
        f.write("""
<!DOCTYPE html>
<html>
<head>
    <title>Redirecting...</title>
    <meta http-equiv="refresh" content="0;URL='/static/frontend.html'" />
</head>
<body>
    <p>Redirecting to the application...</p>
</body>
</html>
        """)
    
    # Start the server
    app.run(debug=True)